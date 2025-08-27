import streamlit as st
from pptx import Presentation
from io import BytesIO

st.set_page_config(page_title="Presentify Demo", layout="wide")
st.title("Presentify — Demo & Submission-ready Viewer")
st.write("Upload a `.pptx` template to preview slides, and (optionally) test generation by sending text to the backend API.")

uploaded = st.file_uploader("Upload .pptx/.potx template", type=["pptx","potx"])

if uploaded:
    try:
        prs = Presentation(uploaded)
    except Exception as e:
        st.error("Failed to read presentation: " + str(e))
        st.stop()

    slides = prs.slides if prs else []
    total = len(slides)
    if total == 0:
        st.warning("No slides found in uploaded file.")
    else:
        slide_idx = st.slider("Slide", 1, total, 1)
        st.markdown(f"**Slide {slide_idx} / {total}**")
        slide = slides[slide_idx - 1]

        cols = st.columns(2)
        with cols[0]:
            st.subheader("Text content")
            for shape in slide.shapes:
                try:
                    if getattr(shape, 'has_text_frame', False) and shape.text.strip():
                        st.write(shape.text.strip())
                except Exception:
                    pass
        with cols[1]:
            st.subheader("Images (if any)")
            for shape in slide.shapes:
                try:
                    if shape.shape_type == 13:  # picture
                        st.image(shape.image.blob, use_column_width=True)
                except Exception:
                    pass

        st.progress(slide_idx/total)

st.markdown("---")
st.subheader("Quick generation demo (optional)")
text = st.text_area("Sample text to generate slides from", height=160)
provider = st.selectbox("LLM provider (local test)", ["openai","anthropic","gemini"])
api_key = st.text_input("API key (used only for demo)", type="password")
template_file = st.file_uploader("Template for generation (optional)", type=["pptx","potx"], key="gen_template")
if st.button("Generate (demo)"):
    if not text.strip():
        st.warning("Please provide sample text for generation.")
    elif not api_key.strip():
        st.warning("Please provide your provider API key for the demo.")
    elif not template_file:
        st.warning("Please upload a template to generate into.")
    else:
        import requests
        try:
            files = {'template': (template_file.name, template_file.getvalue(), 'application/vnd.openxmlformats-officedocument.presentationml.presentation')}
            data = {'text': text, 'provider': provider, 'api_key': api_key}
            resp = requests.post("http://localhost:8000/api/generate", data=data, files=files, timeout=120)
            if resp.status_code == 200:
                st.success("Generated presentation received — click to download.")
                st.download_button("Download Generated Deck", resp.content, file_name="generated_demo.pptx", mime="application/vnd.openxmlformats-officedocument.presentationml.presentation")
            else:
                st.error(f"Generation failed: {resp.status_code} — {resp.text[:200]}")
        except Exception as e:
            st.error("Could not reach local backend at http://localhost:8000. Start backend and try again. Error: " + str(e))