from typing import List, Dict, Any, Optional
from pptx import Presentation
from pptx.util import Inches
from pptx.enum.shapes import MSO_SHAPE_TYPE
from io import BytesIO

def collect_template_images(prs: Presentation) -> List[bytes]:
    imgs = []
    for s in prs.slides:
        for shp in s.shapes:
            if shp.shape_type == MSO_SHAPE_TYPE.PICTURE:
                try:
                    imgs.append(shp.image.blob)
                except Exception:
                    pass
    for layout in prs.slide_layouts:
        for shp in layout.shapes:
            if shp.shape_type == MSO_SHAPE_TYPE.PICTURE:
                try:
                    imgs.append(shp.image.blob)
                except Exception:
                    pass
    return imgs

def find_title_content_layout(prs: Presentation):
    best = None
    for layout in prs.slide_layouts:
        try:
            types = [getattr(ph.placeholder_format, 'type', None) for ph in layout.placeholders]
            has_title = any(t == 1 for t in types if t is not None)
            has_body = any(t == 2 for t in types if t is not None)
            if has_title and has_body:
                return layout
            if has_title:
                best = best or layout
        except Exception:
            continue
    return best or prs.slide_layouts[0]

def apply_slide(slide, item: Dict[str, Any], reuse_imgs: List[bytes], img_idx: int) -> int:
    title = item.get("title", "")[:120]
    bullets = [str(b)[:300] for b in item.get("bullets", [])][:8]
    notes = item.get("notes")

    try:
        if slide.shapes.title:
            slide.shapes.title.text = title
    except Exception:
        tx = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(1.2))
        tx.text_frame.text = title

    body = None
    for ph in slide.placeholders:
        try:
            if ph.placeholder_format.type == 2:
                body = ph.text_frame
                break
        except Exception:
            continue

    if body is None:
        body = slide.shapes.add_textbox(Inches(0.7), Inches(1.6), Inches(8.5), Inches(4)).text_frame

    try:
        body.clear()
    except Exception:
        pass

    for i, b in enumerate(bullets):
        if i == 0:
            body.text = b
        else:
            p = body.add_paragraph()
            p.text = b
            p.level = 0

    if reuse_imgs and (img_idx % 2 == 0):
        try:
            slide.shapes.add_picture(BytesIO(reuse_imgs[img_idx % len(reuse_imgs)]), Inches(5), Inches(1.6), width=Inches(4))
            img_idx += 1
        except Exception:
            pass

    if notes:
        try:
            slide.notes_slide.notes_text_frame.text = notes
        except Exception:
            pass

    return img_idx

def build_presentation(template_blob: bytes, outline: List[Dict[str, Any]], title_fallback: Optional[str]="Presentation") -> bytes:
    tmpl_io = BytesIO(template_blob)
    prs = Presentation(tmpl_io)
    layout = find_title_content_layout(prs)
    imgs = collect_template_images(prs)
    img_idx = 0

    for item in outline:
        slide = prs.slides.add_slide(layout)
        img_idx = apply_slide(slide, item, imgs, img_idx)

    try:
        if prs.slides and prs.slides[0].shapes.title and not prs.slides[0].shapes.title.text.strip():
            prs.slides[0].shapes.title.text = title_fallback
    except Exception:
        pass

    out = BytesIO()
    prs.save(out)
    return out.getvalue()